import csv
from pptx import Presentation
import gpt

def cards_from_ppt(ppt_file, csv_file):
    ppt = Presentation(ppt_file)
    count_slides = 1
    # Iterate over each slide
    # Open the CSV file in write mode
    with open(csv_file, 'w', newline='') as csvfile:
        csvwriter = csv.writer(csvfile)
        total_slides = len(ppt.slides)
        # Iterate over each pair of slides
        for i in range(0, total_slides, 2):
            print(f"\nSlides {i+1} and {i+2}")
            slide_text = ""
            # Iterate over the current slide and the next slide
            for j in range(i, min(i+2, total_slides)):
                for shape in ppt.slides[j].shapes:
                    if shape.has_text_frame:
                        for paragraph in shape.text_frame.paragraphs:
                            for run in paragraph.runs:
                                # Add the text to the slide_text variable
                                slide_text += run.text
            # remove school appearances in the text so it doesn't affect the  prompt
            slide_text = slide_text.replace("Grand Canyon University", "")
            # Call the gpt function
            print(slide_text)
            response = gpt.call(slide_text)
            lines = response.split('\n')
            # Initialize variables to hold question and answer
            question = None
            answer = None

            # Iterate over each line in the response
            for line in lines:
                print(line)
                # Clean data
                if "Question:" in line or "(Question:" in line:
                    question = line.replace("(Question:", "").replace("Question:", "").strip()
                elif "Answer:" in line or "(Answer:" in line:
                    answer = line.replace("(Answer:", "").replace("Answer:", "").strip()

                # If both question and answer are found, write them to the CSV file
                if question and answer:
                    print("Adding flashcard to CSV")
                    csvwriter.writerow([question, answer])
                    question = None
                    answer = None

#get user's inout for the ppt file path
path = input("Enter the path to the ppt file: ")
ppt_file = path # path to the ppt file
csv_file = 'output.csv'
cards_from_ppt(ppt_file, csv_file)
